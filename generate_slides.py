from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9 Widescreen standard
prs.slide_height = Inches(7.5)

# Styling constants
DARK_BLUE = RGBColor(11, 28, 62)
TEXT_GRAY = RGBColor(60, 60, 60)
LIGHT_BG = RGBColor(245, 247, 250)

def add_title(slide, text):
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.33), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = 'Arial'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    return title_box

# ==========================================
# SLIDE 1: Title Slide (Dark Background)
# ==========================================
blank_layout = prs.slide_layouts[6]
slide1 = prs.slides.add_slide(blank_layout)

# Background color
bg = slide1.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BLUE
bg.line.fill.background()

# Title text
title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(2.0))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Private AI Model Benchmarking: Telco Blueprint"
p.font.name = 'Arial'
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

p2 = tf.add_paragraph()
p2.text = "Comparing Qwen3-30B MoE vs. Qwen2.5-7B Dense on Cloudera AI"
p2.font.name = 'Arial'
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(180, 200, 220)
p2.space_before = Pt(20)


# ==========================================
# SLIDE 2: Infrastructure Configuration
# ==========================================
slide2 = prs.slides.add_slide(blank_layout)
add_title(slide2, "Infrastructure & Cost Configuration")

# Add a table for the specs
rows, cols = 6, 3
left, top, width, height = Inches(0.5), Inches(1.5), Inches(12.33), Inches(4.0)
table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height)
table = table_shape.table

# Column widths
table.columns[0].width = Inches(3.33)
table.columns[1].width = Inches(4.5)
table.columns[2].width = Inches(4.5)

headers = ["Metric", "Model A (The Heavyweight)", "Model B (The Lightweight)"]
data = [
    ["Model", "Qwen3-Coder-30B-A3B-Instruct", "Qwen2.5-Coder-7B-Instruct"],
    ["AWS Instance", "g5.12xlarge", "g5.2xlarge"],
    ["Compute / Memory", "4x A10 GPUs, 16 CPUs, 96GiB RAM", "1x A10 GPU, 6 CPUs, 26GiB RAM"],
    ["Cost per Hour", "$5.672 / hr", "$1.212 / hr"],
    ["Hardware Scaling", "Provides 4x the GPU resources", "Baseline profile"]
]

for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE
    cell.text_frame.paragraphs[0].font.bold = True
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

for row_idx, row_data in enumerate(data):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        if col_idx == 0:
            cell.text_frame.paragraphs[0].font.bold = True

# Context note text box
note_box = slide2.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.33), Inches(1.0))
tf = note_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Cost Analysis: The g5.12xlarge infrastructure scales directly with hardware capabilities, running roughly 4.6x more expensive than the g5.2xlarge footprint."
p.font.name = 'Arial'
p.font.size = Pt(14)
p.font.color.rgb = TEXT_GRAY
p.font.italic = True


# ==========================================
# SLIDE 3: Latency & Total Time (Placeholder Layout)
# ==========================================
slide3 = prs.slides.add_slide(blank_layout)
add_title(slide3, "Model Performance & Latency Metrics")

content_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(5.0))
tf = content_box.text_frame
tf.word_wrap = True

def add_bullet(tf, text, level=0, bold=False):
    p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
    p.text = text
    p.level = level
    p.font.name = 'Arial'
    p.font.size = Pt(18) if level == 0 else Pt(16)
    p.font.bold = bold
    p.font.color.rgb = DARK_BLUE if level == 0 else TEXT_GRAY
    p.space_after = Pt(10)

add_bullet(tf, "Qwen3-Coder-30B-A3B-Instruct (4x A10 GPUs)", bold=True)
add_bullet(tf, "• Time-to-First-Token (TTFT): [Insert your value] ms", level=1)
add_bullet(tf, "• Workflow Execution Duration: [Insert total time] mins", level=1)
add_bullet(tf, "• MoE Efficiency: Utilizes sparse Mixture-of-Experts routing to minimize inference latency.", level=1)

add_bullet(tf, "Qwen2.5-Coder-7B-Instruct (1x A10 GPU)", bold=True)
add_bullet(tf, "• Time-to-First-Token (TTFT): [Insert your value] ms", level=1)
add_bullet(tf, "• Workflow Execution Duration: [Insert total time] mins", level=1)
add_bullet(tf, "• Compute Constraint: Every token forces an explicit pass over all 7 billion parameters.", level=1)


# ==========================================
# SLIDE 4: Model Decisions & Recommendations
# ==========================================
slide4 = prs.slides.add_slide(blank_layout)
add_title(slide4, "Agentic Quality & Recommendations")

content_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(5.0))
tf4 = content_box4.text_frame
tf4.word_wrap = True

add_bullet(tf4, "Autonomous Decision Making Quality", bold=True)
add_bullet(tf4, "• Qwen3-30B: Exhibited highly stable schema adherence and multi-step orchestration capability.", level=1)
add_bullet(tf4, "• Qwen2.5-7B: Highly effective for routine execution but more susceptible to schema formatting drops.", level=1)

add_bullet(tf4, "Strategic Deployment Options", bold=True)
add_bullet(tf4, "• Strategic Move A (Autonomy): Deploy Qwen3-30B for dynamic, zero-touch network orchestration workflows.", level=1)
add_bullet(tf4, "• Strategic Move B (Cost Efficiency): Implement a hybrid 'Sidekick' routing pattern—routing 80% of routine processing to the 7B instance ($1.21/hr) and utilizing the 30B instance ($5.67/hr) as an escalation layer.", level=1)

# Save the presentation
prs.save("Telco_Model_Benchmarking.pptx")
print("Successfully generated 'Telco_Model_Benchmarking.pptx'!")